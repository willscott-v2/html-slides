"""
build-pptx.py
Build a PPTX from slide screenshots + speaker notes.
Run after export-screenshots.js has generated the screenshots folder + notes.json.

Requirements:
    pip install python-pptx

Usage:
    python3 build-pptx.py
    python3 build-pptx.py --input path/to/screenshots/ --output my-talk.pptx

Output: a 16:9 Widescreen PPTX at 13.333" × 7.5" with full-bleed 1920×1080
images and speaker notes pulled from notes.json (written by
export-screenshots.js).
"""

import argparse
import glob
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


parser = argparse.ArgumentParser(description="Build PPTX from slide screenshots")
parser.add_argument("--input",  default=os.path.join(os.path.dirname(__file__), "template", "screenshots"),
                    help="Folder containing slide_01.png, slide_02.png, ... and notes.json")
parser.add_argument("--output", default=os.path.join(os.path.dirname(__file__), "presentation.pptx"),
                    help="Output .pptx file path")
args = parser.parse_args()

SLIDE_DIR = args.input
OUTPUT    = args.output

if not os.path.isdir(SLIDE_DIR):
    print(f"Error: screenshot folder not found: {SLIDE_DIR}")
    print("Run export-screenshots.js first.")
    sys.exit(1)

screenshots = sorted(glob.glob(os.path.join(SLIDE_DIR, "slide_*.png")))
if not screenshots:
    print(f"Error: no slide_*.png files found in: {SLIDE_DIR}")
    sys.exit(1)

notes_path = os.path.join(SLIDE_DIR, "notes.json")
if os.path.isfile(notes_path):
    with open(notes_path, encoding="utf-8") as f:
        notes = json.load(f)
    print(f"Loaded {len(notes)} speaker notes from {os.path.basename(notes_path)}")
else:
    notes = []
    print(f"Warning: {notes_path} not found — slides will have no speaker notes.")
    print("         (Re-run export-screenshots.js to generate it.)")

print(f"Building PPTX from {len(screenshots)} screenshots...")
print(f"Output: {OUTPUT}\n")

# 16:9 Widescreen — matches PowerPoint / Keynote default, 1920×1080 fills edge-to-edge
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

for i, img_path in enumerate(screenshots):
    slide = prs.slides.add_slide(blank_layout)

    slide.shapes.add_picture(
        img_path,
        left=Inches(0), top=Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )

    if i < len(notes) and notes[i]:
        slide.notes_slide.notes_text_frame.text = notes[i]

    print(f"  ✓ Slide {i + 1} / {len(screenshots)}")

prs.save(OUTPUT)
print(f"\nSaved → {OUTPUT}")
